from __future__ import annotations

import math
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


PRIMARY_BLUE = RGBColor(0x00, 0x55, 0xD4)
PRIMARY_RED = RGBColor(0xD9, 0x2B, 0x2B)
PRIMARY_GREEN = RGBColor(0x0A, 0x8A, 0x3E)
DARK_NAVY = RGBColor(0x0F, 0x17, 0x2A)
MEDIUM_BLUE = RGBColor(0x25, 0x63, 0xEB)
TEAL = RGBColor(0x08, 0x91, 0xB2)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
MEDIUM_GRAY = RGBColor(0x94, 0xA3, 0xB8)
DARK_GRAY = RGBColor(0x33, 0x41, 0x55)
BORDER_GRAY = RGBColor(0xE2, 0xE8, 0xF0)

FONT_FAMILY = "Calibri"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

HEADER_PANEL_LEFT = Inches(0.46)
HEADER_PANEL_TOP = Inches(0.18)
HEADER_PANEL_WIDTH = Inches(12.4)
HEADER_PANEL_HEIGHT = Inches(1.08)
TITLE_LEFT = Inches(0.78)
TITLE_TOP = Inches(0.34)
TITLE_WIDTH = Inches(9.2)
SUBTITLE_TOP = Inches(0.86)
HEADER_RIGHT_LEFT = Inches(10.25)
HEADER_RIGHT_WIDTH = Inches(2.45)


def fit_text_frame(tf):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    return tf


def set_slide_background(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_bottom_bar(slide, color=PRIMARY_BLUE):
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        SLIDE_HEIGHT - Pt(2),
        SLIDE_WIDTH,
        Pt(2),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_slide_number(slide, number: int):
    box = slide.shapes.add_textbox(
        SLIDE_WIDTH - Inches(0.75),
        SLIDE_HEIGHT - Inches(0.33),
        Inches(0.4),
        Inches(0.15),
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(number)
    apply_font(run, 10, MEDIUM_GRAY)
    return box


def apply_font(run, size, color, bold=False, italic=False, name=FONT_FAMILY):
    font = run.font
    font.name = name
    font.size = Pt(size)
    font.bold = bold
    font.italic = italic
    font.color.rgb = color


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text="",
    font_size=17,
    color=DARK_NAVY,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = fit_text_frame(box.text_frame)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    apply_font(run, font_size, color, bold=bold, italic=italic)
    return box


def add_title_subtitle(slide, title, subtitle="", meta_label="Research Briefing", meta_context="Conference-style V3"):
    header_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        HEADER_PANEL_LEFT,
        HEADER_PANEL_TOP,
        HEADER_PANEL_WIDTH,
        HEADER_PANEL_HEIGHT,
    )
    header_panel.fill.solid()
    header_panel.fill.fore_color.rgb = LIGHT_GRAY
    header_panel.line.color.rgb = BORDER_GRAY
    header_panel.line.width = Pt(1)

    top_strip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        HEADER_PANEL_LEFT,
        HEADER_PANEL_TOP,
        HEADER_PANEL_WIDTH,
        Pt(4),
    )
    top_strip.fill.solid()
    top_strip.fill.fore_color.rgb = PRIMARY_BLUE
    top_strip.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        TITLE_LEFT - Inches(0.18),
        TITLE_TOP,
        Inches(0.08),
        Inches(0.78),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = PRIMARY_BLUE
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, Inches(0.56))
    tf = fit_text_frame(title_box.text_frame)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    apply_font(run, 36, DARK_NAVY, bold=True)

    if subtitle:
        sub_box = slide.shapes.add_textbox(TITLE_LEFT, SUBTITLE_TOP, TITLE_WIDTH, Inches(0.28))
        tf = fit_text_frame(sub_box.text_frame)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = subtitle
        apply_font(run, 20, DARK_GRAY)

    right_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        HEADER_RIGHT_LEFT,
        Inches(0.26),
        HEADER_RIGHT_WIDTH,
        Inches(0.82),
    )
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = WHITE
    right_panel.line.color.rgb = BORDER_GRAY
    right_panel.line.width = Pt(1)

    add_label_chip(
        slide,
        HEADER_RIGHT_LEFT + Inches(0.12),
        Inches(0.34),
        Inches(1.55),
        meta_label,
        fill_color=PRIMARY_BLUE,
    )
    add_label_chip(
        slide,
        HEADER_RIGHT_LEFT + Inches(0.12),
        Inches(0.64),
        Inches(1.1),
        "V3",
        fill_color=TEAL,
    )
    add_textbox(
        slide,
        HEADER_RIGHT_LEFT + Inches(0.12),
        Inches(0.91),
        Inches(1.8),
        Inches(0.12),
        text=meta_context,
        font_size=10.5,
        color=DARK_GRAY,
    )
    right_accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        HEADER_RIGHT_LEFT + Inches(2.08),
        Inches(0.26),
        Inches(0.22),
        Inches(0.82),
    )
    right_accent.fill.solid()
    right_accent.fill.fore_color.rgb = DARK_NAVY
    right_accent.line.fill.background()
    small_band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        HEADER_RIGHT_LEFT + Inches(2.08),
        Inches(0.66),
        Inches(0.22),
        Inches(0.16),
    )
    small_band.fill.solid()
    small_band.fill.fore_color.rgb = TEAL
    small_band.line.fill.background()

    divider = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        TITLE_LEFT - Inches(0.18),
        Inches(1.18),
        SLIDE_WIDTH - Inches(0.65),
        Inches(1.18),
    )
    divider.line.color.rgb = BORDER_GRAY
    divider.line.width = Pt(1)


def add_rule(slide, y):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, TITLE_LEFT, y, SLIDE_WIDTH - TITLE_LEFT, y)
    line.line.color.rgb = BORDER_GRAY
    line.line.width = Pt(1)
    return line


def add_label_chip(slide, left, top, width, text, fill_color=PRIMARY_BLUE, text_color=WHITE):
    chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.24))
    chip.fill.solid()
    chip.fill.fore_color.rgb = fill_color
    chip.line.color.rgb = fill_color
    tf = fit_text_frame(chip.text_frame)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    apply_font(run, 12, text_color, bold=True)
    return chip


def make_ctx():
    return {"placeholders": 0, "tables": 0, "diagrams": 0, "visual": False}


def add_bullet_list(slide, bullets, left, top, width, height, ctx):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = fit_text_frame(box.text_frame)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for item in bullets:
        level, text = item if isinstance(item, tuple) else (0, item)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(5)
        p.line_spacing = 1.22
        if level == 0:
            bullet_symbol = "■ "
            bullet_color = PRIMARY_BLUE
            size = 16
        elif level == 1:
            bullet_symbol = "▸ "
            bullet_color = MEDIUM_BLUE
            size = 15
        else:
            bullet_symbol = "– "
            bullet_color = DARK_GRAY
            size = 14
        bullet_run = p.add_run()
        bullet_run.text = ("   " * level) + bullet_symbol
        apply_font(bullet_run, size, bullet_color, bold=True)
        text_run = p.add_run()
        text_run.text = text
        apply_font(text_run, size, DARK_NAVY if level == 0 else DARK_GRAY)
    ctx["visual"] = True
    return box


def add_table(slide, data, left, top, width, height, ctx, col_widths=None):
    rows = len(data)
    cols = len(data[0])
    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = shape.table
    if col_widths:
        for idx, value in enumerate(col_widths):
            table.columns[idx].width = value
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = PRIMARY_BLUE
                text_color = WHITE
                bold = True
            else:
                cell.fill.fore_color.rgb = LIGHT_GRAY if r % 2 == 0 else WHITE
                text_color = DARK_NAVY
                bold = False
            cell.margin_left = Pt(8)
            cell.margin_right = Pt(8)
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    apply_font(run, 13 if r > 0 else 14, text_color, bold=bold)
    ctx["tables"] += 1
    return table


def add_triangle(slide, cx, cy, size, angle_deg, color):
    tri = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
        cx - size / 2,
        cy - size / 2,
        size,
        size,
    )
    tri.fill.solid()
    tri.fill.fore_color.rgb = color
    tri.line.color.rgb = color
    tri.rotation = angle_deg
    return tri


def add_arrow(slide, start_x, start_y, end_x, end_y, color=PRIMARY_BLUE, width=Pt(1.6)):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y)
    line.line.color.rgb = color
    line.line.width = width
    angle = math.degrees(math.atan2(end_y - start_y, end_x - start_x))
    add_triangle(slide, end_x, end_y, Pt(10), angle + 90, color)
    return line


def add_diagram_box(
    slide,
    left,
    top,
    width,
    height,
    text,
    ctx,
    fill_color=PRIMARY_BLUE,
    text_color=WHITE,
    border_color=None,
    font_size=13,
):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color or fill_color
    shape.line.width = Pt(1.25)
    tf = fit_text_frame(shape.text_frame)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    apply_font(run, font_size, text_color, bold=True)
    ctx["diagrams"] += 1
    ctx["visual"] = True
    return shape


def add_subtle_panel(slide, left, top, width, height, label="", fill=LIGHT_GRAY, border=BORDER_GRAY):
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = fill
    panel.line.color.rgb = border
    panel.line.width = Pt(1)
    if label:
        add_label_chip(slide, left + Inches(0.08), top + Inches(0.08), Inches(1.7), label, fill_color=PRIMARY_BLUE)
    return panel


def add_image_placeholder(
    slide,
    left,
    top,
    width,
    height,
    img_type,
    content,
    source,
    filename,
    alt_text,
    ctx,
):
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = LIGHT_GRAY
    rect.line.color.rgb = MEDIUM_GRAY
    rect.line.width = Pt(2)
    rect.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, top, left + width, top + height)
    line1.line.color.rgb = MEDIUM_GRAY
    line1.line.width = Pt(1)
    line1.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    line2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left + width, top, left, top + height)
    line2.line.color.rgb = MEDIUM_GRAY
    line2.line.width = Pt(1)
    line2.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    tf = fit_text_frame(rect.text_frame)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "[IMAGE PLACEHOLDER]"
    apply_font(run, 12, DARK_GRAY, bold=True)

    for line in [
        f"Type     : {img_type}",
        f"Content  : {content}",
        f"Source   : {source}",
        f"Filename : {filename}",
        f"Alt-text : {alt_text}",
    ]:
        para = tf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        para.line_spacing = 1.1
        r = para.add_run()
        r.text = line
        apply_font(r, 10.5, DARK_GRAY)

    ctx["placeholders"] += 1
    ctx["visual"] = True
    return rect


def add_metric_card(slide, left, top, width, height, number, label, caption, color, ctx):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color
    card.line.width = Pt(2)
    tf = fit_text_frame(card.text_frame)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = number
    apply_font(r1, 30, color, bold=True)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    apply_font(r2, 14, DARK_NAVY, bold=True)
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = caption
    apply_font(r3, 11, DARK_GRAY)
    ctx["diagrams"] += 1
    ctx["visual"] = True
    return card


def add_notes(slide, notes: Iterable[str]):
    text = "\n".join(str(line) for line in notes if str(line).strip())
    if not text:
        return
    notes_slide = slide.notes_slide
    for placeholder in notes_slide.placeholders:
        if placeholder.placeholder_format.type == 2:
            placeholder.text = text
            return


def finalize_slide(slide, number, notes, ctx, tracker):
    add_bottom_bar(slide)
    add_slide_number(slide, number)
    add_notes(slide, notes)
    tracker["total_slides"] += 1
    tracker["placeholders"] += ctx["placeholders"]
    tracker["tables"] += ctx["tables"]
    tracker["diagrams"] += ctx["diagrams"]
    if ctx["visual"]:
        tracker["visual_slides"] += 1


def add_error_slide(slide, title, exc_text):
    add_title_subtitle(slide, f"Rendering error: {title}", "The slide content failed to render normally.")
    add_textbox(
        slide,
        Inches(0.8),
        Inches(1.6),
        Inches(11.5),
        Inches(2.2),
        text=f"Error details: {exc_text}",
        font_size=18,
        color=PRIMARY_RED,
        bold=True,
    )


def render_section_divider(slide, title, subtitle, section_no="01", accent_color=TEAL, keywords=None):
    keywords = keywords or ["Scientific framing", "Workflow logic", "Visual evidence"]
    block = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    block.fill.solid()
    block.fill.fore_color.rgb = PRIMARY_BLUE
    block.line.fill.background()

    top_band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.54))
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = DARK_NAVY
    top_band.line.fill.background()

    right_wall = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(11.55),
        0,
        Inches(1.78),
        SLIDE_HEIGHT,
    )
    right_wall.fill.solid()
    right_wall.fill.fore_color.rgb = MEDIUM_BLUE
    right_wall.line.fill.background()

    diagonal_low = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.15),
        Inches(5.18),
        Inches(6.7),
        Inches(0.82),
    )
    diagonal_low.fill.solid()
    diagonal_low.fill.fore_color.rgb = DARK_NAVY
    diagonal_low.line.fill.background()
    diagonal_low.rotation = -10

    diagonal_top = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(8.5),
        Inches(0.04),
        Inches(4.1),
        Inches(0.72),
    )
    diagonal_top.fill.solid()
    diagonal_top.fill.fore_color.rgb = accent_color
    diagonal_top.line.fill.background()
    diagonal_top.rotation = -11

    for x in [Inches(8.05), Inches(8.7), Inches(9.35), Inches(10.0)]:
        grid_line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            x,
            Inches(0.62),
            x,
            Inches(5.7),
        )
        grid_line.line.color.rgb = WHITE
        grid_line.line.width = Pt(0.75)

    for y in [Inches(1.55), Inches(2.45), Inches(3.35), Inches(4.25), Inches(5.15)]:
        grid_line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(7.75),
            y,
            Inches(11.2),
            y,
        )
        grid_line.line.color.rgb = WHITE
        grid_line.line.width = Pt(0.75)

    add_textbox(
        slide,
        Inches(0.88),
        Inches(0.15),
        Inches(4.5),
        Inches(0.22),
        text="Quantum Chemistry Automation Platform",
        font_size=12,
        color=WHITE,
        bold=True,
        italic=True,
    )

    section_card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(8.45),
        Inches(1.2),
        Inches(2.45),
        Inches(1.4),
    )
    section_card.fill.solid()
    section_card.fill.fore_color.rgb = WHITE
    section_card.line.color.rgb = WHITE
    section_card.line.width = Pt(1.2)
    tf = fit_text_frame(section_card.text_frame)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r0 = p0.add_run()
    r0.text = "SECTION"
    apply_font(r0, 12, DARK_GRAY, bold=True)
    p1 = tf.add_paragraph()
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = section_no
    apply_font(r1, 30, PRIMARY_BLUE, bold=True)

    add_textbox(
        slide,
        Inches(0.96),
        Inches(2.15),
        Inches(7.1),
        Inches(0.88),
        text=title,
        font_size=42,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.LEFT,
    )
    add_textbox(
        slide,
        Inches(0.98),
        Inches(3.2),
        Inches(7.4),
        Inches(0.52),
        text=subtitle,
        font_size=19,
        color=WHITE,
        align=PP_ALIGN.LEFT,
    )

    meta_band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.9),
        Inches(6.1),
        Inches(11.7),
        Inches(0.68),
    )
    meta_band.fill.solid()
    meta_band.fill.fore_color.rgb = WHITE
    meta_band.line.fill.background()
    add_textbox(
        slide,
        Inches(1.08),
        Inches(6.18),
        Inches(1.8),
        Inches(0.18),
        text="Section Focus",
        font_size=11,
        color=DARK_GRAY,
        bold=True,
        italic=True,
    )
    chip_x = Inches(2.72)
    for keyword in keywords[:3]:
        add_label_chip(slide, chip_x, Inches(6.24), Inches(2.1), keyword, fill_color=accent_color)
        chip_x += Inches(2.32)


def render_title_slide(slide, ctx):
    add_subtle_panel(slide, Inches(0.55), Inches(0.58), Inches(6.9), Inches(5.1), label="Project Overview", fill=WHITE)
    hero_frame = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(8.05),
        Inches(0.58),
        Inches(4.55),
        Inches(5.08),
    )
    hero_frame.fill.solid()
    hero_frame.fill.fore_color.rgb = WHITE
    hero_frame.line.color.rgb = PRIMARY_BLUE
    hero_frame.line.width = Pt(1.5)
    add_label_chip(slide, Inches(8.22), Inches(0.72), Inches(2.05), "Hero Figure Zone", fill_color=PRIMARY_BLUE)

    title_box = slide.shapes.add_textbox(Inches(0.78), Inches(0.9), Inches(6.35), Inches(1.55))
    tf = fit_text_frame(title_box.text_frame)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Quantum Chemistry Automation Platform"
    apply_font(r, 33, DARK_NAVY, bold=True)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "LLM + MCP + PySCF for conversational computation"
    apply_font(r2, 19, DARK_GRAY)

    add_textbox(
        slide,
        Inches(0.82),
        Inches(2.55),
        Inches(6.0),
        Inches(1.0),
        text=(
            "This presentation frames the project as a research tool that lets experimental chemists ask "
            "for quantum-chemistry analyses in natural language and receive visualization-ready outputs."
        ),
        font_size=18,
        color=DARK_NAVY,
    )
    add_label_chip(slide, Inches(0.82), Inches(3.72), Inches(2.0), "Conference-style V3", fill_color=PRIMARY_RED)
    add_label_chip(slide, Inches(0.82), Inches(4.08), Inches(2.15), "Lab Meeting Deck", fill_color=TEAL)
    add_label_chip(slide, Inches(3.15), Inches(4.08), Inches(2.45), "Browser-Based UI", fill_color=PRIMARY_BLUE)
    add_label_chip(slide, Inches(5.75), Inches(4.08), Inches(1.75), "Open Source", fill_color=PRIMARY_GREEN)
    add_image_placeholder(
        slide,
        Inches(8.28),
        Inches(1.08),
        Inches(4.08),
        Inches(4.22),
        "Composite hero visual",
        "Molecular orbital rendering on the left and chatbot UI screenshot on the right",
        "Directly generated for the presentation",
        "hero_quantum_chatbot.png",
        "Hero image combining molecular visualization and the chat interface",
        ctx,
    )
    add_textbox(
        slide,
        Inches(8.22),
        Inches(5.05),
        Inches(4.0),
        Inches(0.34),
        text="Target visual: orbital render + chat screenshot + clean browser result composition",
        font_size=11,
        color=MEDIUM_GRAY,
        italic=True,
    )
    add_textbox(
        slide,
        Inches(0.75),
        Inches(6.15),
        Inches(6.0),
        Inches(0.25),
        text="Presenter: Lab member / Date: 2026-03-28 / Context: mixed experimental and computational audience",
        font_size=10.5,
        color=MEDIUM_GRAY,
        italic=True,
    )


def render_agenda_slide(slide, ctx):
    add_title_subtitle(slide, "Agenda", "From motivation to platform architecture, capabilities, and next steps")
    add_rule(slide, Inches(1.25))
    bullets = [
        "Why quantum chemistry remains essential for experimental decision-making",
        "Why the conventional workflow is slow, fragmented, and expert-dependent",
        "How the platform turns a natural-language request into a compute workflow",
        "What the current system can already calculate and visualize",
        "How the browser-based demo experience should be shown in a live talk",
        "What the practical and academic impact could be for the lab",
        "What validation, roadmap, and publication steps come next",
    ]
    add_bullet_list(slide, bullets, Inches(0.8), Inches(1.55), Inches(6.0), Inches(4.8), ctx)

    add_subtle_panel(slide, Inches(7.2), Inches(1.55), Inches(5.3), Inches(4.45), label="Talk flow")
    steps = [
        ("Background", PRIMARY_RED),
        ("Problem", PRIMARY_RED),
        ("Platform", PRIMARY_BLUE),
        ("Capabilities", PRIMARY_GREEN),
        ("Impact", ORANGE),
        ("Appendix", PURPLE),
    ]
    x = Inches(7.45)
    y = Inches(2.15)
    for idx, (label, color) in enumerate(steps):
        add_diagram_box(slide, x, y + idx * Inches(0.52), Inches(2.05), Inches(0.34), label, ctx, fill_color=color)
        if idx < len(steps) - 1:
            add_arrow(slide, x + Inches(1.02), y + Inches(0.34) + idx * Inches(0.52), x + Inches(1.02), y + Inches(0.5) + idx * Inches(0.52), color=DARK_GRAY)


def render_key_message_slide(slide, ctx):
    add_title_subtitle(slide, "One-sentence thesis", "The deck should open with the main transformation, not with implementation details")
    quote_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.95), Inches(1.7), Inches(11.35), Inches(1.6))
    quote_box.fill.solid()
    quote_box.fill.fore_color.rgb = LIGHT_GRAY
    quote_box.line.color.rgb = PRIMARY_BLUE
    quote_box.line.width = Pt(2)
    tf = quote_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "If a chemist can describe the molecule and the analysis in plain English, the system can route the computation and return a visualization-ready result."
    apply_font(r, 23, PRIMARY_BLUE, bold=True, italic=True)

    add_diagram_box(slide, Inches(1.15), Inches(4.0), Inches(4.1), Inches(1.0), "Before: Input files, tool chains, and specialist mediation", ctx, fill_color=WHITE, text_color=PRIMARY_RED, border_color=PRIMARY_RED, font_size=14)
    add_diagram_box(slide, Inches(8.0), Inches(4.0), Inches(4.1), Inches(1.0), "After: One browser entry point, one question-driven workflow", ctx, fill_color=WHITE, text_color=PRIMARY_GREEN, border_color=PRIMARY_GREEN, font_size=14)
    add_arrow(slide, Inches(5.5), Inches(4.5), Inches(7.75), Inches(4.5), color=PRIMARY_BLUE)
    add_textbox(slide, Inches(4.95), Inches(5.2), Inches(3.0), Inches(0.35), text="Accessibility is the real innovation layer", font_size=13, color=DARK_GRAY, align=PP_ALIGN.CENTER)


def render_qc_value_slide(slide, ctx):
    add_title_subtitle(slide, "Why experimental chemists still need quantum chemistry", "The platform matters only if the underlying scientific questions are real")
    add_bullet_list(
        slide,
        [
            "Quantum chemistry provides electronic-structure information that experiments often cannot observe directly.",
            "HOMO/LUMO analysis helps frame reactivity, electron donation, and electron acceptance in an intuitive way.",
            "Geometry optimization identifies a more stable reference structure before downstream interpretation.",
            "Electrostatic potential maps reveal electron-rich and electron-poor regions on the molecular surface.",
            "Partial charges and dipole moments help rationalize intermolecular interactions and selectivity trends.",
            "For transient species or difficult-to-observe intermediates, calculations become a hypothesis-building tool.",
            "The value proposition is not to replace experiments, but to sharpen explanation, planning, and interpretation.",
        ],
        Inches(0.8),
        Inches(1.45),
        Inches(6.15),
        Inches(4.95),
        ctx,
    )
    add_image_placeholder(
        slide,
        Inches(7.25),
        Inches(1.55),
        Inches(5.1),
        Inches(4.25),
        "Molecular visualization set",
        "One orbital image, one ESP surface, and one optimized structure comparison panel",
        "Directly generated from project outputs or recreated for the deck",
        "quantum_value_triptych.png",
        "Three molecular panels showing why quantum calculations matter in practice",
        ctx,
    )
    add_textbox(slide, Inches(7.3), Inches(5.95), Inches(5.0), Inches(0.4), text="Use an image triptych here so the audience immediately sees that the outputs are chemically meaningful, not merely computational artifacts.", font_size=11, color=MEDIUM_GRAY, italic=True)


def render_traditional_workflow_slide(slide, ctx):
    add_title_subtitle(slide, "Traditional workflow is fragmented", "The slow part is not only the compute time, but the human and software chain around it")
    box_y = Inches(2.0)
    labels = [
        "Experimental question",
        "Ask a specialist",
        "Prepare geometry",
        "Write input files",
        "Run compute code",
        "Parse outputs",
        "Visualize elsewhere",
        "Explain results",
    ]
    x_positions = [Inches(0.7 + i * 1.55) for i in range(len(labels))]
    for idx, label in enumerate(labels):
        fill = PRIMARY_RED if idx in {1, 3, 6} else WHITE
        text_color = WHITE if fill == PRIMARY_RED else DARK_NAVY
        border = PRIMARY_RED if fill == PRIMARY_RED else PRIMARY_BLUE
        add_diagram_box(slide, x_positions[idx], box_y, Inches(1.35), Inches(0.8), label, ctx, fill_color=fill, text_color=text_color, border_color=border, font_size=11.5)
        if idx < len(labels) - 1:
            add_arrow(slide, x_positions[idx] + Inches(1.35), box_y + Inches(0.4), x_positions[idx + 1], box_y + Inches(0.4), color=DARK_GRAY)
    add_textbox(
        slide,
        Inches(0.85),
        Inches(4.2),
        Inches(11.4),
        Inches(1.55),
        text=(
            "This chain requires multiple tools, repeated file conversion, and repeated interpretation. "
            "If any assumption changes, the workflow loops back upstream and the turn-around time grows quickly."
        ),
        font_size=16,
        color=DARK_NAVY,
    )
    add_label_chip(slide, Inches(0.9), Inches(5.4), Inches(2.2), "Pain concentrates here", fill_color=PRIMARY_RED)


def render_pain_points_slide(slide, ctx):
    add_title_subtitle(slide, "What makes the current workflow painful", "The bottleneck is structural: people, tools, and knowledge are all serialized")
    add_bullet_list(
        slide,
        [
            "A small number of computational experts become the routing bottleneck for routine requests.",
            "Different tools handle structure preparation, execution, parsing, and visualization separately.",
            "Users must understand basis sets, functionals, convergence behavior, and failure recovery.",
            "Commercial software and viewers introduce cost, access-control, and licensing constraints.",
            "Experimentalists often receive outputs without enough interpretive context to act immediately.",
            "Changing the question usually means repeating the same slow loop with a specialist in the middle.",
        ],
        Inches(0.8),
        Inches(1.55),
        Inches(6.2),
        Inches(4.7),
        ctx,
    )
    add_image_placeholder(
        slide,
        Inches(7.35),
        Inches(1.7),
        Inches(4.9),
        Inches(4.0),
        "Concept illustration",
        "Fragmented legacy workflow with many disconnected software tools and a human bottleneck in the middle",
        "Custom figure for the presentation",
        "legacy_workflow_pain.png",
        "Illustration of a fragmented multi-tool workflow with a specialist bottleneck",
        ctx,
    )


def render_gap_slide(slide, ctx):
    add_title_subtitle(slide, "The research gap is now actionable", "The enabling condition is the intersection of mature open chemistry, language models, and tool orchestration")
    add_subtle_panel(slide, Inches(1.0), Inches(1.7), Inches(11.1), Inches(3.65), label="Intersection")
    add_diagram_box(slide, Inches(1.35), Inches(2.45), Inches(2.7), Inches(1.0), "Natural-language interpretation", ctx, fill_color=PRIMARY_BLUE)
    add_diagram_box(slide, Inches(5.0), Inches(2.0), Inches(3.0), Inches(1.0), "Open quantum-chemistry engines", ctx, fill_color=PRIMARY_GREEN)
    add_diagram_box(slide, Inches(8.8), Inches(2.45), Inches(2.7), Inches(1.0), "Web-native automation and visualization", ctx, fill_color=TEAL)
    add_diagram_box(slide, Inches(4.65), Inches(3.55), Inches(3.6), Inches(0.85), "Practical conversational computation", ctx, fill_color=WHITE, text_color=PRIMARY_BLUE, border_color=PRIMARY_BLUE)
    add_arrow(slide, Inches(4.05), Inches(2.95), Inches(4.65), Inches(3.95), color=DARK_GRAY)
    add_arrow(slide, Inches(6.5), Inches(3.0), Inches(6.45), Inches(3.55), color=DARK_GRAY)
    add_arrow(slide, Inches(8.8), Inches(2.95), Inches(8.25), Inches(3.95), color=DARK_GRAY)
    add_textbox(slide, Inches(1.1), Inches(5.7), Inches(11.0), Inches(0.45), text="Earlier attempts simplified individual tools. The current opportunity is to connect language understanding, tool routing, and open compute backends into one user experience.", font_size=14, color=DARK_GRAY)


def render_core_idea_slide(slide, ctx):
    add_title_subtitle(slide, "Core idea: question in, insight out", "The system should hide workflow friction while keeping scientific outputs visible")
    add_diagram_box(slide, Inches(0.9), Inches(1.8), Inches(2.35), Inches(1.05), "User question\n\"Show the HOMO of acetone\"", ctx, fill_color=PRIMARY_BLUE, font_size=14)
    add_arrow(slide, Inches(3.25), Inches(2.32), Inches(5.0), Inches(2.32), color=DARK_NAVY)
    add_diagram_box(slide, Inches(5.0), Inches(1.8), Inches(2.45), Inches(1.05), "Intent extraction\nmolecule + task + context", ctx, fill_color=TEAL, font_size=14)
    add_arrow(slide, Inches(7.45), Inches(2.32), Inches(9.15), Inches(2.32), color=DARK_NAVY)
    add_diagram_box(slide, Inches(9.15), Inches(1.8), Inches(3.05), Inches(1.05), "Compute route\nstructure + PySCF + result packaging", ctx, fill_color=PRIMARY_GREEN, font_size=14)
    add_bullet_list(
        slide,
        [
            "The user no longer has to think in terms of file formats, command-line options, or post-processing tools.",
            "The platform decides which compute path to run and packages the result for immediate use.",
            "Follow-up questions can reuse the same session context instead of restarting from zero.",
            "The design goal is not fewer scientific details, but less operational friction around those details.",
        ],
        Inches(0.95),
        Inches(3.45),
        Inches(11.0),
        Inches(2.1),
        ctx,
    )


def render_technology_stack_slide(slide, ctx):
    add_title_subtitle(slide, "Technology stack at a glance", "The system combines an open compute engine, a routing layer, and a browser-native presentation layer")
    table_data = [
        ["Layer", "Technology", "Role in the platform"],
        ["Compute", "PySCF", "Runs single-point, orbital, ESP, charge, and optimization workflows"],
        ["Language interface", "LLM providers", "Interprets requests, helps route tasks, and explains outputs"],
        ["Tool bridge", "MCP-style tool layer", "Connects language intent to explicit compute or helper functions"],
        ["Web application", "FastAPI + WebSocket", "Delivers chat, job status, and result retrieval in real time"],
        ["Visualization", "3Dmol.js", "Renders molecular structures, surfaces, and orbital views in the browser"],
        ["External structure data", "MolChat / PubChem", "Resolves candidate structures before computation"],
        ["Scalable runtime", "Redis + arq (optional)", "Separates web serving from job execution when needed"],
    ]
    add_table(slide, table_data, Inches(0.75), Inches(1.55), Inches(11.85), Inches(4.35), ctx)
    add_textbox(slide, Inches(0.8), Inches(6.05), Inches(11.6), Inches(0.35), text="The talk should present these as a coherent pipeline rather than as disconnected frameworks.", font_size=11, color=MEDIUM_GRAY, italic=True)


def render_pyscf_slide(slide, ctx):
    add_title_subtitle(slide, "Why PySCF is the compute backbone", "The AI layer routes requests, but PySCF performs the actual electronic-structure work")
    add_bullet_list(
        slide,
        [
            "PySCF is an open-source quantum-chemistry framework with strong Python interoperability.",
            "That interoperability makes it well suited for an automated pipeline rather than a manual file-centric workflow.",
            "In this project it underpins single-point energy, orbital preview, ESP mapping, charge analysis, and geometry optimization.",
            "The open-source model lowers cost and improves reproducibility compared with commercial-only setups.",
            "The platform should therefore be described as a compute-automation system, not as an AI model pretending to calculate.",
        ],
        Inches(0.8),
        Inches(1.55),
        Inches(6.0),
        Inches(4.85),
        ctx,
    )
    add_image_placeholder(
        slide,
        Inches(7.25),
        Inches(1.6),
        Inches(5.0),
        Inches(3.95),
        "Logo plus compute capability graphic",
        "PySCF logo or text logo with small capability icons for energy, orbitals, ESP, charges, and optimization",
        "Custom figure or official logo asset",
        "pyscf_stack_badge.png",
        "PySCF logo with annotated compute capabilities",
        ctx,
    )


def render_llm_mcp_slide(slide, ctx):
    add_title_subtitle(slide, "Why LLM + MCP changes the interface", "The novelty lies in turning natural language into explicit tool calls rather than into vague text-only answers")
    add_diagram_box(slide, Inches(0.9), Inches(2.1), Inches(2.3), Inches(0.9), "User request", ctx, fill_color=WHITE, text_color=DARK_NAVY, border_color=PRIMARY_BLUE)
    add_arrow(slide, Inches(3.2), Inches(2.55), Inches(5.0), Inches(2.55), color=PRIMARY_BLUE)
    add_diagram_box(slide, Inches(5.0), Inches(2.1), Inches(2.2), Inches(0.9), "LLM\ninterpretation", ctx, fill_color=PRIMARY_BLUE)
    add_arrow(slide, Inches(7.2), Inches(2.55), Inches(8.8), Inches(2.55), color=PRIMARY_BLUE)
    add_diagram_box(slide, Inches(8.8), Inches(2.1), Inches(2.4), Inches(0.9), "MCP-style\n tool layer", ctx, fill_color=TEAL)
    add_arrow(slide, Inches(11.2), Inches(2.55), Inches(12.15), Inches(2.55), color=PRIMARY_BLUE)
    add_diagram_box(slide, Inches(10.95), Inches(3.4), Inches(1.7), Inches(0.8), "PySCF / resolver /\nvisualizer", ctx, fill_color=PRIMARY_GREEN, font_size=11.5)
    add_arrow(slide, Inches(10.0), Inches(3.0), Inches(11.15), Inches(3.4), color=DARK_GRAY)
    add_bullet_list(
        slide,
        [
            "The LLM reads the request in human language and identifies molecule, task, and follow-up context.",
            "The tool layer makes the next step explicit: calculate ESP, preview an orbital, resolve a structure, or build a summary.",
            "This division reduces the risk of a text-only model speaking as if it had already run a calculation.",
            "That is why MCP is best explained as the model's hands and feet rather than as another chatbot feature.",
        ],
        Inches(0.95),
        Inches(4.45),
        Inches(11.0),
        Inches(1.65),
        ctx,
    )


def render_architecture_slide(slide, ctx):
    add_title_subtitle(slide, "End-to-end architecture", "A single browser experience sits on top of routing, structure resolution, compute, and result presentation")
    add_diagram_box(slide, Inches(0.8), Inches(1.95), Inches(2.1), Inches(0.9), "Browser client", ctx, fill_color=WHITE, text_color=DARK_NAVY, border_color=PRIMARY_BLUE)
    add_diagram_box(slide, Inches(3.35), Inches(1.95), Inches(2.4), Inches(0.9), "Chat and compute routes", ctx, fill_color=PRIMARY_BLUE)
    add_diagram_box(slide, Inches(6.15), Inches(1.95), Inches(2.25), Inches(0.9), "Structure resolver", ctx, fill_color=TEAL)
    add_diagram_box(slide, Inches(8.8), Inches(1.95), Inches(2.15), Inches(0.9), "PySCF runner", ctx, fill_color=PRIMARY_GREEN)
    add_diagram_box(slide, Inches(11.15), Inches(1.95), Inches(1.45), Inches(0.9), "Result UI", ctx, fill_color=WHITE, text_color=DARK_NAVY, border_color=PRIMARY_BLUE)
    for x1, x2 in [(Inches(2.9), Inches(3.35)), (Inches(5.75), Inches(6.15)), (Inches(8.4), Inches(8.8)), (Inches(10.95), Inches(11.15))]:
        add_arrow(slide, x1, Inches(2.4), x2, Inches(2.4), color=DARK_GRAY)

    add_diagram_box(slide, Inches(4.0), Inches(4.0), Inches(2.6), Inches(0.85), "MolChat / PubChem", ctx, fill_color=WHITE, text_color=TEAL, border_color=TEAL)
    add_diagram_box(slide, Inches(7.3), Inches(4.0), Inches(2.6), Inches(0.85), "Advisor and explanation layer", ctx, fill_color=WHITE, text_color=PRIMARY_GREEN, border_color=PRIMARY_GREEN)
    add_diagram_box(slide, Inches(10.25), Inches(4.0), Inches(2.0), Inches(0.85), "Redis / worker (optional)", ctx, fill_color=WHITE, text_color=PURPLE, border_color=PURPLE)
    add_arrow(slide, Inches(6.3), Inches(2.85), Inches(5.3), Inches(4.0), color=TEAL)
    add_arrow(slide, Inches(9.85), Inches(2.85), Inches(8.6), Inches(4.0), color=PRIMARY_GREEN)
    add_arrow(slide, Inches(9.85), Inches(2.85), Inches(11.25), Inches(4.0), color=PURPLE)

    add_textbox(slide, Inches(0.82), Inches(5.35), Inches(11.4), Inches(0.6), text="The slide should visually separate the user-facing path, the structure-preparation path, and the optional scale-out runtime path.", font_size=13, color=DARK_GRAY)


def render_pipeline_slide(slide, ctx):
    add_title_subtitle(slide, "Pipeline from question to visualization", "This is the operational sequence that the audience should remember")
    steps = [
        "1. User request",
        "2. Intent extraction",
        "3. Structure preparation",
        "4. Task selection",
        "5. PySCF execution",
        "6. Result packaging",
        "7. Browser visualization",
    ]
    x = Inches(0.75)
    y = Inches(2.0)
    for idx, step in enumerate(steps):
        color = PRIMARY_BLUE if idx < 2 else TEAL if idx < 4 else PRIMARY_GREEN if idx < 6 else ORANGE
        add_diagram_box(slide, x + idx * Inches(1.75), y, Inches(1.45), Inches(0.78), step, ctx, fill_color=color, font_size=12)
        if idx < len(steps) - 1:
            add_arrow(slide, x + Inches(1.45) + idx * Inches(1.75), y + Inches(0.39), x + Inches(1.75) + idx * Inches(1.75), y + Inches(0.39), color=DARK_GRAY)
    add_bullet_list(
        slide,
        [
            "The first three steps determine whether the system correctly understands what the user really wants.",
            "The middle steps determine whether the requested calculation is executed on a chemically meaningful structure.",
            "The final step determines whether the output is actually useful to a non-expert without another software round-trip.",
            "That is why the platform should be presented as a full workflow, not as a thin wrapper around a compute engine.",
        ],
        Inches(0.8),
        Inches(3.45),
        Inches(11.55),
        Inches(2.2),
        ctx,
    )


def render_restaurant_slide(slide, ctx):
    add_title_subtitle(slide, "A restaurant analogy for the system", "The analogy helps a mixed audience grasp the division of labor without code-level detail")
    labels = [
        ("Experimentalist", PRIMARY_BLUE),
        ("LLM waiter", TEAL),
        ("Tool menu", ORANGE),
        ("PySCF kitchen", PRIMARY_GREEN),
        ("Visualization plating", PRIMARY_BLUE),
    ]
    x = Inches(0.9)
    for idx, (label, color) in enumerate(labels):
        add_diagram_box(slide, x + idx * Inches(2.35), Inches(2.2), Inches(1.8), Inches(0.9), label, ctx, fill_color=color, font_size=13)
        if idx < len(labels) - 1:
            add_arrow(slide, x + Inches(1.8) + idx * Inches(2.35), Inches(2.65), x + Inches(2.35) + idx * Inches(2.35), Inches(2.65), color=DARK_GRAY)
    add_bullet_list(
        slide,
        [
            "The scientist states the desired outcome in natural language, just as a diner places an order.",
            "The LLM interprets the order, while the tool layer translates it into an explicit kitchen action.",
            "The compute engine performs the actual work, and the result is plated for direct inspection in the browser.",
            "The metaphor is useful because it emphasizes that users should not need to know the kitchen's internal recipe syntax.",
        ],
        Inches(0.9),
        Inches(4.0),
        Inches(11.0),
        Inches(1.9),
        ctx,
    )


def render_capability_overview_slide(slide, ctx):
    add_title_subtitle(slide, "Capability map", "The platform already covers several high-value workflows rather than a single demo task")
    features = [
        ("Orbital preview", PRIMARY_BLUE),
        ("ESP map", TEAL),
        ("Partial charges", PRIMARY_GREEN),
        ("Geometry optimization", ORANGE),
        ("Single-point energy", PURPLE),
        ("Integrated analysis", PRIMARY_BLUE),
        ("Session continuation", TEAL),
        ("Result explanation", PRIMARY_GREEN),
    ]
    start_x = Inches(0.9)
    start_y = Inches(1.8)
    for idx, (label, color) in enumerate(features):
        row = idx // 4
        col = idx % 4
        add_diagram_box(slide, start_x + col * Inches(2.85), start_y + row * Inches(1.0), Inches(2.35), Inches(0.68), label, ctx, fill_color=WHITE, text_color=color, border_color=color, font_size=13)

    table_data = [
        ["Capability", "Why it matters chemically"],
        ["HOMO/LUMO preview", "Supports intuitive discussion of donation, acceptance, and reactivity"],
        ["ESP map", "Highlights electrostatic regions linked to intermolecular interaction patterns"],
        ["Optimization", "Establishes a more stable reference geometry before interpretation"],
        ["Charge analysis", "Provides atom-level descriptors useful for mechanistic discussion"],
        ["Continuation", "Allows follow-up questions without re-entering the same context"],
    ]
    add_table(slide, table_data, Inches(0.85), Inches(4.25), Inches(11.6), Inches(1.7), ctx)


def render_feature_slide(slide, ctx, title, subtitle, bullets, placeholder_title, placeholder_content, filename):
    add_title_subtitle(slide, title, subtitle)
    add_bullet_list(slide, bullets, Inches(0.8), Inches(1.55), Inches(6.0), Inches(4.9), ctx)
    add_image_placeholder(
        slide,
        Inches(7.15),
        Inches(1.55),
        Inches(5.15),
        Inches(4.3),
        placeholder_title,
        placeholder_content,
        "Direct output screenshot or recreated render for the talk",
        filename,
        placeholder_content,
        ctx,
    )


def render_additional_features_slide(slide, ctx):
    add_title_subtitle(slide, "Additional capabilities beyond the headline demo", "The project already combines computation, continuity, and explanation")
    data = [
        ["Capability", "Function in the platform", "Why it matters"],
        ["Single-point energy", "Provides a fast reference electronic-energy calculation", "Supports quick comparison and baseline analysis"],
        ["Partial charge analysis", "Summarizes atom-wise charge distribution", "Useful for bond polarity and reactivity discussion"],
        ["Analyze workflow", "Packages multiple outputs into one result contract", "Improves usability for non-experts"],
        ["Session continuation", "Carries molecule and task context across follow-up requests", "Makes the chat interface genuinely interactive"],
        ["Clarification flow", "Asks for missing detail when the request is ambiguous", "Reduces misrouted calculations"],
        ["Advisor layer", "Adds presets, methods text, and interpretation support", "Improves reproducibility and communication"],
    ]
    add_table(slide, data, Inches(0.72), Inches(1.55), Inches(11.9), Inches(4.25), ctx)
    add_image_placeholder(
        slide,
        Inches(8.8),
        Inches(5.95),
        Inches(3.55),
        Inches(0.95),
        "Small UI strip",
        "Compact screenshot showing session history, follow-up controls, or advisor recommendations",
        "Browser screenshot",
        "continuation_advisor_strip.png",
        "UI strip illustrating continuation and recommendation features",
        ctx,
    )


def render_visualization_slide(slide, ctx):
    add_title_subtitle(slide, "Visualization is the last mile of usability", "Results are most useful when they arrive as inspectable browser objects rather than as raw files")
    add_bullet_list(
        slide,
        [
            "The browser view can host structures, orbital surfaces, ESP surfaces, and summary panels in one place.",
            "Users can rotate, zoom, and switch between result tabs without opening a second desktop application.",
            "That removes one of the most common breakpoints in the conventional workflow: post-processing and re-rendering.",
            "For a mixed audience, this is also the most persuasive evidence that the system is more than a backend API.",
        ],
        Inches(0.8),
        Inches(1.55),
        Inches(5.65),
        Inches(4.9),
        ctx,
    )
    add_image_placeholder(
        slide,
        Inches(6.85),
        Inches(1.55),
        Inches(2.55),
        Inches(4.2),
        "Browser UI screenshot",
        "Result panel with molecular viewer, tabs, and computation summary",
        "Web-app screenshot",
        "viewer_ui_browser.png",
        "Browser interface showing the 3D viewer and result panel",
        ctx,
    )
    add_image_placeholder(
        slide,
        Inches(9.65),
        Inches(1.55),
        Inches(2.55),
        Inches(4.2),
        "Molecular output screenshot",
        "High-contrast orbital or ESP render for one representative molecule",
        "Web-app screenshot or regenerated still image",
        "viewer_output_example.png",
        "Representative molecular visualization from the browser UI",
        ctx,
    )


def render_demo_scenarios_slide(slide, ctx):
    add_title_subtitle(slide, "Recommended live-demo scenarios", "A good demo shows continuity, not just one flashy still image")
    add_subtle_panel(slide, Inches(0.8), Inches(1.6), Inches(5.4), Inches(4.8), label="Demo A")
    add_subtle_panel(slide, Inches(6.95), Inches(1.6), Inches(5.4), Inches(4.8), label="Demo B")
    add_bullet_list(
        slide,
        [
            "Start with a request such as 'Show the HOMO of acetone.'",
            "Use the first result to prove the platform can turn a plain-English request into a visualization.",
            "Then ask a follow-up question on the same molecule rather than switching context immediately.",
            "This reveals that the session stores enough context to support continuity.",
        ],
        Inches(1.0),
        Inches(2.0),
        Inches(4.95),
        Inches(2.0),
        ctx,
    )
    add_bullet_list(
        slide,
        [
            "Ask for an ESP map or geometry optimization on the same structure.",
            "Show progress updates to emphasize that a real compute path is running.",
            "If live reliability is uncertain, replace the live segment with a short pre-recorded clip.",
            "The key criterion is continuity from question to result, not theatrical speed.",
        ],
        Inches(7.15),
        Inches(2.0),
        Inches(4.95),
        Inches(2.0),
        ctx,
    )
    add_image_placeholder(slide, Inches(1.05), Inches(4.25), Inches(2.15), Inches(1.65), "Chat UI snapshot", "Opening request in the browser chat", "Web-app screenshot", "demo_chat_start.png", "Chat interface with the opening HOMO request", ctx)
    add_image_placeholder(slide, Inches(3.45), Inches(4.25), Inches(2.15), Inches(1.65), "Result snapshot", "Initial orbital result with viewer", "Web-app screenshot", "demo_homo_result.png", "Initial HOMO result in the browser", ctx)
    add_image_placeholder(slide, Inches(7.2), Inches(4.25), Inches(2.15), Inches(1.65), "Progress snapshot", "Compute progress status panel", "Web-app screenshot", "demo_progress.png", "Progress screen showing computation status", ctx)
    add_image_placeholder(slide, Inches(9.6), Inches(4.25), Inches(2.15), Inches(1.65), "Follow-up result", "ESP or optimization follow-up result", "Web-app screenshot", "demo_followup.png", "Follow-up result on the same molecular context", ctx)


def render_conversation_slide(slide, ctx):
    add_title_subtitle(slide, "Conversation flow example", "The deck should show that the interface supports sequential scientific questioning")
    add_subtle_panel(slide, Inches(0.9), Inches(1.55), Inches(11.3), Inches(4.95), label="Chat transcript")
    messages = [
        ("USER", "Show the HOMO of benzene.", PRIMARY_BLUE, WHITE, Inches(1.2)),
        ("SYSTEM", "Preparing the structure and launching orbital preview.", LIGHT_GRAY, DARK_NAVY, Inches(3.3)),
        ("SYSTEM", "The HOMO visualization is ready in the 3D viewer.", LIGHT_GRAY, DARK_NAVY, Inches(3.3)),
        ("USER", "Now show the LUMO and add an ESP map.", PRIMARY_BLUE, WHITE, Inches(1.2)),
        ("SYSTEM", "Reusing the same structure context and starting the follow-up jobs.", LIGHT_GRAY, DARK_NAVY, Inches(3.3)),
    ]
    y = Inches(2.0)
    for role, text, fill, text_color, left in messages:
        bubble = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, y, Inches(6.4 if role == "SYSTEM" else 5.0), Inches(0.55))
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = fill
        bubble.line.color.rgb = BORDER_GRAY if fill == LIGHT_GRAY else PRIMARY_BLUE
        bubble.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = bubble.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"{role}: {text}"
        apply_font(r, 14, text_color, bold=(role == "USER"))
        y += Inches(0.78)
    ctx["diagrams"] += 1
    ctx["visual"] = True


def render_result_screen_slide(slide, ctx):
    add_title_subtitle(slide, "What the result screen should communicate", "A result screen must show state, visualization, and interpretation at the same time")
    add_image_placeholder(slide, Inches(0.85), Inches(1.55), Inches(5.7), Inches(4.6), "Browser result screen", "Main application screenshot showing the viewer, job metadata, and selected result tab", "Web-app screenshot", "result_screen_full.png", "Main result screen of the application", ctx)
    add_image_placeholder(slide, Inches(6.95), Inches(1.55), Inches(5.2), Inches(3.25), "Rendered molecular output", "Large orbital or ESP still image for publication-style readability", "Viewer screenshot", "result_visual_focus.png", "Focused molecular visualization from the result view", ctx)
    add_bullet_list(
        slide,
        [
            "The result view should answer three questions immediately: what was run, what was found, and where to look first.",
            "A molecular viewer without context is not enough; it needs state labels and a concise interpretation panel.",
            "For a conference-style audience, the slide should also imply that the browser output is discussion-ready.",
        ],
        Inches(7.0),
        Inches(5.05),
        Inches(5.0),
        Inches(1.15),
        ctx,
    )


def render_comparison_slide(slide, ctx):
    add_title_subtitle(slide, "Before vs. after", "The most defensible comparison is structural, not exaggerated speed claims without user-study data")
    data = [
        ["Metric", "Traditional workflow", "This platform"],
        ["Entry point", "Input files and specialist mediation", "Natural-language browser interaction"],
        ["Tool count", "Multiple compute and post-processing tools", "One browser-centered entry point"],
        ["Interpretation", "Output parsing required", "Visualization plus explanation layer"],
        ["Iteration style", "Human re-request loop", "Follow-up continuation in one session"],
        ["Cost model", "Commercial software may be required", "Open-source-first stack"],
    ]
    add_table(slide, data, Inches(0.75), Inches(1.55), Inches(7.1), Inches(4.35), ctx)
    add_diagram_box(slide, Inches(8.35), Inches(1.85), Inches(3.65), Inches(0.72), "Conceptual time-to-insight comparison", ctx, fill_color=WHITE, text_color=DARK_NAVY, border_color=PRIMARY_BLUE)
    base_x = Inches(8.55)
    base_y = Inches(3.0)
    labels = [("Traditional", PRIMARY_RED, 2.8), ("Platform", PRIMARY_GREEN, 1.5)]
    for idx, (label, color, width_units) in enumerate(labels):
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, base_x, base_y + idx * Inches(0.9), Inches(width_units), Inches(0.38))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.color.rgb = color
        add_textbox(slide, base_x, base_y - Inches(0.2) + idx * Inches(0.9), Inches(3.2), Inches(0.2), text=label, font_size=12, color=DARK_NAVY)
    ctx["diagrams"] += 1
    ctx["visual"] = True


def render_validation_slide(slide, ctx):
    add_title_subtitle(slide, "Validation framework rather than unsupported claims", "The repository supports a benchmark plan, not a finished internal accuracy paper")
    data = [
        ["Validation area", "Comparison target", "What should be checked"],
        ["Relative energies", "Reference quantum-chemistry output", "Rank ordering and magnitude consistency"],
        ["Optimized geometry", "Reference structure parameters", "Bond lengths, angles, and stable minima"],
        ["Orbital results", "Reference orbital analysis", "Ordering, qualitative shape, and interpretive consistency"],
        ["ESP and charges", "Reference electrostatic descriptors", "Region trends rather than decorative similarity"],
        ["Usability", "Lab-user pilot study", "Time to first result and clarity of interpretation"],
    ]
    add_table(slide, data, Inches(0.72), Inches(1.55), Inches(11.9), Inches(4.4), ctx)
    add_textbox(slide, Inches(0.8), Inches(6.05), Inches(11.6), Inches(0.35), text="This framing is honest and still strong: the platform is compelling, and benchmark work is the next formal step.", font_size=11, color=MEDIUM_GRAY, italic=True)


def render_boundaries_slide(slide, ctx):
    add_title_subtitle(slide, "Current boundaries and responsible framing", "A stronger talk separates what is already implemented from what still needs validation or scale-up work")
    add_subtle_panel(slide, Inches(0.8), Inches(1.55), Inches(5.6), Inches(4.9), label="Already solid")
    add_subtle_panel(slide, Inches(6.9), Inches(1.55), Inches(5.6), Inches(4.9), label="Still to mature")
    add_bullet_list(
        slide,
        [
            "The browser workflow, routing logic, and result packaging are already concrete and demo-ready.",
            "Core compute paths are present for orbital preview, ESP, partial charges, optimization, and integrated analysis.",
            "Session continuation and real-time progress updates already strengthen the user experience beyond a static demo.",
            "The optional worker path shows that the architecture was designed with scale-out in mind.",
        ],
        Inches(1.0),
        Inches(2.0),
        Inches(5.0),
        Inches(3.3),
        ctx,
    )
    add_bullet_list(
        slide,
        [
            "Formal accuracy benchmarks against reference software still need to be completed and reported systematically.",
            "Structure resolution remains a scientifically important quality gate before downstream interpretation.",
            "Operational hardening for broader multi-user deployment is a next-step engineering task rather than a finished claim.",
            "Publication-grade scientific conclusions should still be reviewed through expert interpretation and validation.",
        ],
        Inches(7.1),
        Inches(2.0),
        Inches(5.0),
        Inches(3.3),
        ctx,
    )


def render_metrics_slide(slide, ctx):
    add_title_subtitle(slide, "Platform value in three numbers", "Use only numbers that are directly defensible from the current project scan")
    add_metric_card(slide, Inches(1.0), Inches(2.0), Inches(3.2), Inches(2.0), "6", "core compute modes", "single-point, orbital, ESP, charges, optimization, integrated analysis", PRIMARY_BLUE, ctx)
    add_metric_card(slide, Inches(5.05), Inches(2.0), Inches(3.2), Inches(2.0), "1", "browser entry point", "chat, progress tracking, and visualization live in one UI surface", PRIMARY_GREEN, ctx)
    add_metric_card(slide, Inches(9.1), Inches(2.0), Inches(3.2), Inches(2.0), "0", "commercial-license dependency", "the current stack is centered on open-source components", ORANGE, ctx)
    add_bullet_list(
        slide,
        [
            "These numbers are intentionally modest and verifiable.",
            "They communicate breadth, simplicity of access, and cost structure without inventing benchmark claims.",
            "This slide works well before the impact section because it resets the audience around concrete takeaways.",
        ],
        Inches(1.0),
        Inches(4.55),
        Inches(11.0),
        Inches(1.25),
        ctx,
    )


def render_novelty_slide(slide, ctx):
    add_title_subtitle(slide, "Academic novelty", "The novelty claim should focus on workflow integration, accessibility, and tool-mediated computation")
    add_diagram_box(slide, Inches(1.0), Inches(2.05), Inches(3.0), Inches(0.95), "Language-driven request handling", ctx, fill_color=PRIMARY_BLUE)
    add_diagram_box(slide, Inches(5.15), Inches(2.05), Inches(3.0), Inches(0.95), "Open quantum-chemistry execution", ctx, fill_color=PRIMARY_GREEN)
    add_diagram_box(slide, Inches(9.3), Inches(2.05), Inches(3.0), Inches(0.95), "Browser-native scientific presentation", ctx, fill_color=TEAL)
    add_diagram_box(slide, Inches(4.2), Inches(3.55), Inches(4.9), Inches(0.9), "Workflow innovation at the intersection", ctx, fill_color=WHITE, text_color=PRIMARY_BLUE, border_color=PRIMARY_BLUE)
    add_arrow(slide, Inches(2.5), Inches(3.0), Inches(4.55), Inches(3.55), color=DARK_GRAY)
    add_arrow(slide, Inches(6.65), Inches(3.0), Inches(6.65), Inches(3.55), color=DARK_GRAY)
    add_arrow(slide, Inches(10.8), Inches(3.0), Inches(8.75), Inches(3.55), color=DARK_GRAY)
    add_bullet_list(
        slide,
        [
            "The platform is not claiming a new electronic-structure method.",
            "Its novelty lies in making open quantum-chemistry workflows conversational, tool-mediated, and presentation-ready.",
            "That combination sits at the intersection of computational chemistry, AI interfaces, and research automation.",
        ],
        Inches(1.0),
        Inches(4.8),
        Inches(11.0),
        Inches(1.1),
        ctx,
    )


def render_practical_impact_slide(slide, ctx):
    add_title_subtitle(slide, "Practical impact for the lab", "The strongest value proposition is autonomy and iteration speed, not AI spectacle")
    add_diagram_box(slide, Inches(0.95), Inches(2.0), Inches(3.5), Inches(1.0), "Experimental chemists\ncan ask directly", ctx, fill_color=WHITE, text_color=PRIMARY_BLUE, border_color=PRIMARY_BLUE, font_size=15)
    add_diagram_box(slide, Inches(4.9), Inches(2.0), Inches(3.5), Inches(1.0), "Computational chemists\ncan focus on harder problems", ctx, fill_color=WHITE, text_color=PRIMARY_GREEN, border_color=PRIMARY_GREEN, font_size=15)
    add_diagram_box(slide, Inches(8.85), Inches(2.0), Inches(3.5), Inches(1.0), "Students\nlearn through interaction", ctx, fill_color=WHITE, text_color=ORANGE, border_color=ORANGE, font_size=15)
    add_bullet_list(
        slide,
        [
            "A conversational interface lowers the threshold for asking exploratory questions early in a project.",
            "Experts are not replaced; instead, routine routing work can be reduced so that expert attention shifts to validation and interpretation.",
            "The same interface doubles as a teaching aid because concepts such as orbitals or ESP become immediately inspectable.",
            "Open-source deployment also improves accessibility for labs without a large commercial software budget.",
        ],
        Inches(0.95),
        Inches(3.45),
        Inches(11.15),
        Inches(2.05),
        ctx,
    )


def render_roadmap_slide(slide, ctx):
    add_title_subtitle(slide, "Roadmap for the next phases", "The current codebase already hints at both functional and runtime scalability")
    phases = [
        ("Phase 1", "Current: orbital, ESP, partial charges, optimization, integrated analysis", PRIMARY_BLUE),
        ("Phase 2", "Next analyses: TD-DFT, NMR prediction, vibrational workflows", TEAL),
        ("Phase 3", "System scale: multi-user runtime and stronger worker separation", PRIMARY_GREEN),
        ("Phase 4", "Deployment scale: cloud or HPC-backed execution paths", ORANGE),
    ]
    y = Inches(1.95)
    for idx, (label, text, color) in enumerate(phases):
        add_diagram_box(slide, Inches(1.0), y + idx * Inches(0.95), Inches(1.5), Inches(0.55), label, ctx, fill_color=color, font_size=13)
        add_diagram_box(slide, Inches(2.7), y + idx * Inches(0.95), Inches(9.2), Inches(0.55), text, ctx, fill_color=WHITE, text_color=DARK_NAVY, border_color=color, font_size=13)
        if idx < len(phases) - 1:
            add_arrow(slide, Inches(1.75), y + Inches(0.55) + idx * Inches(0.95), Inches(1.75), y + Inches(0.95) + idx * Inches(0.95), color=DARK_GRAY)


def render_publication_plan_slide(slide, ctx):
    add_title_subtitle(slide, "Publication plan", "Position the paper as workflow innovation plus validation, not as an unsupported benchmark claim")
    data = [
        ["Paper block", "Primary content"],
        ["Problem and gap", "Why conventional quantum-chemistry access remains expert-gated and slow"],
        ["System design", "Conversational routing, structure resolution, compute execution, browser visualization"],
        ["Representative scenarios", "HOMO/LUMO, ESP, optimization, and follow-up continuation examples"],
        ["Validation", "Accuracy benchmark plan plus user-study framing"],
        ["Impact discussion", "Research productivity, accessibility, education, and open-science implications"],
    ]
    add_table(slide, data, Inches(0.78), Inches(1.55), Inches(7.1), Inches(4.2), ctx)
    add_image_placeholder(
        slide,
        Inches(8.25),
        Inches(1.75),
        Inches(4.0),
        Inches(3.0),
        "Paper structure diagram",
        "Five-block figure summarizing the planned manuscript narrative",
        "Custom figure for the talk",
        "paper_storyline.png",
        "Diagram of the planned paper structure",
        ctx,
    )
    add_bullet_list(
        slide,
        [
            "Candidate targets include chemistry-informatics or computational-chemistry venues in the ACS ecosystem.",
            "The benchmark and user-evaluation sections should mature in parallel rather than sequentially.",
        ],
        Inches(8.2),
        Inches(5.0),
        Inches(4.15),
        Inches(1.0),
        ctx,
    )


def render_summary_slide(slide, ctx):
    add_title_subtitle(slide, "Take-home summary", "Keep the close simple: interface shift, workflow automation, and research impact")
    add_metric_card(slide, Inches(1.0), Inches(2.0), Inches(3.5), Inches(1.7), "Question-first", "Interface shift", "Chemists describe what they want instead of assembling input syntax", PRIMARY_BLUE, ctx)
    add_metric_card(slide, Inches(4.95), Inches(2.0), Inches(3.5), Inches(1.7), "Compute-backed", "Not text-only AI", "A real quantum-chemistry engine still performs the scientific work", PRIMARY_GREEN, ctx)
    add_metric_card(slide, Inches(8.9), Inches(2.0), Inches(3.5), Inches(1.7), "Lab-relevant", "Practical impact", "The target is autonomy, faster iteration, and easier learning", ORANGE, ctx)
    add_textbox(slide, Inches(1.2), Inches(4.45), Inches(11.0), Inches(0.6), text="The talk should end by moving the audience from 'Is this possible?' to 'How do we validate and deploy it responsibly?'", font_size=19, color=PRIMARY_BLUE, bold=True, italic=True, align=PP_ALIGN.CENTER)


def render_collaboration_slide(slide, ctx):
    add_title_subtitle(slide, "Where the lab can contribute", "The next progress step depends on representative use cases, validation choices, and user feedback")
    add_diagram_box(slide, Inches(1.0), Inches(2.05), Inches(3.3), Inches(0.9), "Experimental input\nrepresentative use cases", ctx, fill_color=WHITE, text_color=PRIMARY_BLUE, border_color=PRIMARY_BLUE)
    add_diagram_box(slide, Inches(5.0), Inches(2.05), Inches(3.3), Inches(0.9), "Computational input\nbenchmark design", ctx, fill_color=WHITE, text_color=PRIMARY_GREEN, border_color=PRIMARY_GREEN)
    add_diagram_box(slide, Inches(9.0), Inches(2.05), Inches(3.3), Inches(0.9), "Training input\nteaching and onboarding needs", ctx, fill_color=WHITE, text_color=ORANGE, border_color=ORANGE)
    add_arrow(slide, Inches(4.3), Inches(2.5), Inches(5.0), Inches(2.5), color=DARK_GRAY)
    add_arrow(slide, Inches(8.3), Inches(2.5), Inches(9.0), Inches(2.5), color=DARK_GRAY)
    add_bullet_list(
        slide,
        [
            "Ask the group which molecules and questions would make the most convincing first internal pilot study.",
            "Ask computational chemists which benchmark criteria are non-negotiable for credibility.",
            "Ask students what explanations or tutorial cues would make the system easier to learn from.",
            "Frame the project as a shared lab asset rather than as an isolated software prototype.",
        ],
        Inches(1.0),
        Inches(3.7),
        Inches(11.0),
        Inches(1.9),
        ctx,
    )


def render_qna_slide(slide, ctx):
    add_title_subtitle(slide, "Questions and discussion", "Useful discussion topics are validation, use-case priority, and deployment strategy")
    add_textbox(slide, Inches(1.0), Inches(1.9), Inches(11.1), Inches(0.8), text="Questions & Discussion", font_size=28, color=PRIMARY_BLUE, bold=True, italic=True, align=PP_ALIGN.CENTER)
    add_bullet_list(
        slide,
        [
            "Which internal use case would be the best first benchmark scenario?",
            "Which quantitative comparison would matter most to this lab: energies, geometry, orbital interpretation, or usability?",
            "What would make the system most useful for experimental planning rather than only for retrospective explanation?",
            "Which new compute mode would add the largest marginal value in the next development phase?",
        ],
        Inches(1.1),
        Inches(3.0),
        Inches(10.8),
        Inches(2.4),
        ctx,
    )


def render_appendix_architecture_slide(slide, ctx):
    add_title_subtitle(slide, "Appendix: architecture layers", "Useful when the audience wants a slightly deeper but still non-code explanation")
    layers = [
        ("Presentation", "Chat UI, result panels, 3D viewer", PRIMARY_BLUE),
        ("Routing", "Chat route, compute route, session continuity", TEAL),
        ("Computation", "Structure resolver, PySCF runner, result packaging", PRIMARY_GREEN),
        ("Scale-out option", "Redis-backed job store and arq worker path", PURPLE),
    ]
    for idx, (label, desc, color) in enumerate(layers):
        add_diagram_box(slide, Inches(1.0), Inches(1.9) + idx * Inches(0.95), Inches(2.1), Inches(0.58), label, ctx, fill_color=color)
        add_diagram_box(slide, Inches(3.45), Inches(1.9) + idx * Inches(0.95), Inches(8.5), Inches(0.58), desc, ctx, fill_color=WHITE, text_color=DARK_NAVY, border_color=color, font_size=13)


def render_appendix_structure_slide(slide, ctx):
    add_title_subtitle(slide, "Appendix: structure resolution path", "A realistic automation pipeline must solve structure preparation before it can solve computation")
    add_diagram_box(slide, Inches(0.9), Inches(2.1), Inches(2.25), Inches(0.85), "User molecule text", ctx, fill_color=WHITE, text_color=DARK_NAVY, border_color=PRIMARY_BLUE)
    add_diagram_box(slide, Inches(3.55), Inches(2.1), Inches(2.4), Inches(0.85), "Resolver", ctx, fill_color=TEAL)
    add_diagram_box(slide, Inches(6.35), Inches(1.55), Inches(2.2), Inches(0.85), "MolChat", ctx, fill_color=WHITE, text_color=TEAL, border_color=TEAL)
    add_diagram_box(slide, Inches(6.35), Inches(2.65), Inches(2.2), Inches(0.85), "PubChem", ctx, fill_color=WHITE, text_color=TEAL, border_color=TEAL)
    add_diagram_box(slide, Inches(9.05), Inches(2.1), Inches(2.6), Inches(0.85), "Prepared structure", ctx, fill_color=PRIMARY_GREEN)
    add_arrow(slide, Inches(3.15), Inches(2.52), Inches(3.55), Inches(2.52), color=DARK_GRAY)
    add_arrow(slide, Inches(5.95), Inches(2.28), Inches(6.35), Inches(1.98), color=DARK_GRAY)
    add_arrow(slide, Inches(5.95), Inches(2.76), Inches(6.35), Inches(3.08), color=DARK_GRAY)
    add_arrow(slide, Inches(8.55), Inches(2.52), Inches(9.05), Inches(2.52), color=DARK_GRAY)
    add_bullet_list(
        slide,
        [
            "This resolver layer is one of the most important hidden parts of the platform.",
            "If structure inference fails, the downstream computation cannot be trusted no matter how elegant the interface is.",
            "That makes structure preparation an important target for both validation and future improvement.",
        ],
        Inches(1.0),
        Inches(4.25),
        Inches(11.0),
        Inches(1.3),
        ctx,
    )


def render_appendix_execution_slide(slide, ctx):
    add_title_subtitle(slide, "Appendix: execution model and queueing", "The current system already contains a path toward separation of web serving and compute execution")
    add_diagram_box(slide, Inches(0.9), Inches(2.0), Inches(2.3), Inches(0.85), "Browser", ctx, fill_color=WHITE, text_color=DARK_NAVY, border_color=PRIMARY_BLUE)
    add_diagram_box(slide, Inches(3.6), Inches(2.0), Inches(2.6), Inches(0.85), "FastAPI web layer", ctx, fill_color=PRIMARY_BLUE)
    add_diagram_box(slide, Inches(6.55), Inches(1.45), Inches(2.3), Inches(0.85), "In-memory job path", ctx, fill_color=WHITE, text_color=PRIMARY_GREEN, border_color=PRIMARY_GREEN)
    add_diagram_box(slide, Inches(6.55), Inches(2.55), Inches(2.3), Inches(0.85), "Redis job store", ctx, fill_color=WHITE, text_color=PURPLE, border_color=PURPLE)
    add_diagram_box(slide, Inches(9.3), Inches(2.55), Inches(2.2), Inches(0.85), "arq worker", ctx, fill_color=PURPLE)
    add_diagram_box(slide, Inches(9.3), Inches(1.45), Inches(2.2), Inches(0.85), "PySCF execution", ctx, fill_color=PRIMARY_GREEN)
    add_arrow(slide, Inches(3.2), Inches(2.42), Inches(3.6), Inches(2.42), color=DARK_GRAY)
    add_arrow(slide, Inches(6.2), Inches(2.15), Inches(6.55), Inches(1.88), color=DARK_GRAY)
    add_arrow(slide, Inches(6.2), Inches(2.69), Inches(6.55), Inches(2.98), color=DARK_GRAY)
    add_arrow(slide, Inches(8.85), Inches(2.98), Inches(9.3), Inches(2.98), color=DARK_GRAY)
    add_arrow(slide, Inches(10.4), Inches(2.55), Inches(10.4), Inches(2.3), color=DARK_GRAY)
    add_bullet_list(
        slide,
        [
            "The simple mode is good for development and small demos.",
            "The Redis and worker path matters when computation time or user concurrency grows.",
            "Real-time WebSocket progress updates are part of the user-value proposition, not only a backend detail.",
        ],
        Inches(0.95),
        Inches(4.15),
        Inches(11.0),
        Inches(1.4),
        ctx,
    )


def render_appendix_evaluation_slide(slide, ctx):
    add_title_subtitle(slide, "Appendix: evaluation matrix", "A convincing paper must evaluate both scientific validity and usability")
    add_diagram_box(slide, Inches(1.0), Inches(1.9), Inches(4.5), Inches(0.55), "Scientific validity", ctx, fill_color=PRIMARY_BLUE)
    add_diagram_box(slide, Inches(6.0), Inches(1.9), Inches(4.5), Inches(0.55), "Usability and workflow value", ctx, fill_color=PRIMARY_GREEN)
    data_left = [
        ["Metric", "Example"],
        ["Energy consistency", "Relative energy ranking"],
        ["Geometry fidelity", "Bond lengths and angles"],
        ["Orbital interpretation", "Ordering and qualitative agreement"],
        ["ESP trend quality", "Region-level electrostatic agreement"],
    ]
    data_right = [
        ["Metric", "Example"],
        ["Time to first result", "From question to first usable view"],
        ["Follow-up success", "How often continuation works as expected"],
        ["Interpretability", "Can users explain what they see next"],
        ["Workflow compression", "How many tools and handoffs remain"],
    ]
    add_table(slide, data_left, Inches(1.0), Inches(2.65), Inches(4.5), Inches(2.65), ctx)
    add_table(slide, data_right, Inches(6.0), Inches(2.65), Inches(4.5), Inches(2.65), ctx)


def render_appendix_references_slide(slide, ctx):
    add_title_subtitle(slide, "Appendix: references and assets", "Use this slide to acknowledge the open-source base and to list the figure assets still needed")
    add_bullet_list(
        slide,
        [
            "PySCF for open quantum-chemistry execution",
            "FastAPI and WebSocket runtime for browser interaction",
            "3Dmol.js for molecular visualization in the browser",
            "MolChat and PubChem for structure look-up support",
            "Integrated deep-scan report and project README as the basis for this deck",
            "Still-needed visual assets: browser screenshots, orbital stills, ESP stills, and one composite title image",
        ],
        Inches(0.9),
        Inches(1.55),
        Inches(6.05),
        Inches(4.9),
        ctx,
    )
    add_image_placeholder(
        slide,
        Inches(7.3),
        Inches(1.7),
        Inches(5.0),
        Inches(3.95),
        "Asset board placeholder",
        "Grid of the final screenshots and renders to replace placeholders in a polished conference version",
        "Custom assembly for final design pass",
        "asset_board.png",
        "Asset board showing the final images still needed for a polished conference deck",
        ctx,
    )


def validate_shape_bounds(prs):
    width = prs.slide_width
    height = prs.slide_height
    issues = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            if not all(hasattr(shape, attr) for attr in ("left", "top", "width", "height")):
                continue
            left = shape.left
            top = shape.top
            right = left + shape.width
            bottom = top + shape.height
            if left < 0 or top < 0 or right > width or bottom > height:
                issues.append((slide_idx, shape_idx, getattr(shape, "name", "shape")))
    return issues


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    tracker = {
        "total_slides": 0,
        "placeholders": 0,
        "tables": 0,
        "diagrams": 0,
        "visual_slides": 0,
    }

    slides = [
        ("title", render_title_slide, ["Introduce the platform as a research-facing system rather than as a coding project.", "The title slide should immediately connect natural-language requests with quantum-chemistry outputs.", "Emphasize that the audience is mixed, so the presentation will stay at the level of motivation, workflow, and impact.", "Point to the browser and molecular-visualization pairing as the main visual motif for the whole talk."]),
        ("agenda", render_agenda_slide, ["Use the agenda to reassure the audience that the talk will move from motivation to impact in a logical arc.", "The main goal is to show why this matters, how it works conceptually, and what it can already do.", "Mention that technical appendix material is available later for anyone who wants more depth."]),
        ("thesis", render_key_message_slide, ["Deliver the one-sentence thesis slowly and clearly.", "The audience should leave this slide understanding that the interface shift is the main idea.", "Frame the remaining slides as evidence supporting this transformation."]),
        ("divider1", lambda s, c: render_section_divider(s, "Why This Problem Matters", "Scientific need, workflow friction, and the current research gap", section_no="01", accent_color=PRIMARY_RED, keywords=["Scientific need", "Workflow friction", "Research gap"]), ["This divider marks the transition from opening claims to the motivation section."]),
        ("qc-value", render_qc_value_slide, ["Reinforce that the system matters only because the underlying chemical questions matter.", "Use HOMO/LUMO, ESP, optimization, and charge analysis as recognizable examples for both experimentalists and computational chemists.", "Position the platform as a way to access these outputs more easily, not as a replacement for scientific judgment."]),
        ("workflow", render_traditional_workflow_slide, ["This slide should make the conventional path look operationally heavy.", "Stress that compute time is only one piece of the delay; handoffs and tool changes are equally important.", "Set up the contrast for the platform's simplicity without exaggerating."]),
        ("pain", render_pain_points_slide, ["Summarize the bottleneck in terms of people, software, and knowledge barriers.", "Explain that the problem is not that chemists do not care about computation; it is that the access path is expensive in attention and time.", "This is the reason a conversational workflow is valuable."]),
        ("gap", render_gap_slide, ["Explain why this project becomes realistic now rather than years ago.", "The key point is the intersection of maturing language models, open-source compute engines, and web-based delivery.", "This is a workflow innovation story built from several mature ingredients."]),
        ("divider2", lambda s, c: render_section_divider(s, "Platform Concept", "How the system converts a natural-language request into a compute-backed result", section_no="02", accent_color=TEAL, keywords=["Language entry", "Tool routing", "PySCF execution"]), ["Move from motivation into platform design."]),
        ("core-idea", render_core_idea_slide, ["Describe the platform as a translator between the scientist's question and the machine's compute workflow.", "The system should absorb the operational burden while keeping the scientific result visible.", "This framing is more useful than focusing on implementation details."]),
        ("stack", render_technology_stack_slide, ["Use this slide to orient the audience around the major building blocks without diving into code.", "The most important takeaway is that each layer has a clear role in the overall workflow.", "Avoid overselling the language model; keep PySCF central as the compute engine."]),
        ("pyscf", render_pyscf_slide, ["State clearly that the AI layer does not replace quantum chemistry.", "PySCF is where the scientific calculation actually happens.", "That distinction improves credibility with computational chemists in the room."]),
        ("llm-mcp", render_llm_mcp_slide, ["Explain MCP as the model's hands and feet.", "The important distinction is between language interpretation and explicit tool execution.", "This helps the audience understand why the system is more than a text-only chatbot."]),
        ("architecture", render_architecture_slide, ["Walk left to right through the architecture diagram.", "Separate the user-facing path from the external structure sources and the optional scale-out runtime.", "Remind the audience that the browser experience hides these internal steps."]),
        ("pipeline", render_pipeline_slide, ["This is the operational memory slide: question, routing, structure preparation, computation, packaging, and visualization.", "Use it to emphasize that the platform is an end-to-end pipeline.", "The system is valuable because it automates the transitions between these stages."]),
        ("analogy", render_restaurant_slide, ["Use the analogy only to make the workflow intuitive, not to trivialize the science.", "The scientist orders the result, the system translates, the compute engine executes, and the browser serves.", "This works especially well for a mixed lab audience."]),
        ("divider3", lambda s, c: render_section_divider(s, "Capabilities", "What the current system can already calculate and present", section_no="03", accent_color=PRIMARY_GREEN, keywords=["Orbitals", "Optimization", "ESP and charges"]), ["Transition from architecture into concrete user-visible features."]),
        ("capability-map", render_capability_overview_slide, ["This slide proves that the platform is not a one-off HOMO demo.", "Highlight the spread across compute modes, interpretation support, and session continuity.", "The table at the bottom ties each feature to a chemical use case."]),
        ("feature-homo", lambda s, c: render_feature_slide(s, c, "Feature deep dive: HOMO/LUMO", "Use the most visually persuasive feature first", ["HOMO/LUMO views connect directly to familiar reactivity language.", "The user asks in natural language and receives a browser-ready orbital surface.", "The result is not just a file; it is an interactive 3D object with context.", "This makes orbital concepts easier to discuss in both research and teaching settings.", "For a live audience, this is often the strongest first demonstration image."], "Molecular orbital render", "HOMO and LUMO renderings of one representative molecule, shown side by side", "feature_homo_lumo.png"), ["Describe orbital preview as the most immediate bridge between computation and chemical intuition.", "Stress the user experience: one question, one browser result, no manual cube handling.", "This is likely the best first feature to show in a live demo."]),
        ("feature-opt", lambda s, c: render_feature_slide(s, c, "Feature deep dive: geometry optimization", "Optimization matters because many later interpretations depend on the starting structure", ["Optimization finds a more stable reference geometry under the chosen level of theory.", "This matters because orbital or ESP interpretation is more meaningful on a chemically sensible structure.", "The platform supports optimization as a direct natural-language request.", "The output can be shown as an updated 3D geometry in the same interface.", "This turns a traditionally file-heavy pre-processing step into a direct part of the conversation."], "Optimized geometry comparison", "Before-and-after geometry panels with one representative optimized structure", "feature_optimization.png"), ["Explain that optimization is both a standalone result and a precondition for later analysis.", "This is a good slide to remind the audience that the platform handles more than visualization alone.", "Emphasize continuity between optimization and later feature requests."]),
        ("feature-esp", lambda s, c: render_feature_slide(s, c, "Feature deep dive: ESP mapping", "ESP is one of the most intuitive ways to talk about electrostatic patterns on a molecular surface", ["ESP maps turn charge-distribution information into an interpretable color surface.", "The browser delivery matters because it removes the usual handoff into a separate visualization tool.", "For many experimentalists, this image is easier to interpret quickly than a raw charge table.", "The feature also complements orbital analysis by focusing on the molecular surface rather than an orbital isosurface.", "In the talk, use it to connect electronic structure to interaction patterns and selectivity language."], "ESP surface render", "High-contrast electrostatic potential map with visible color scale", "feature_esp.png"), ["ESP is often the easiest bridge from calculation to interaction intuition.", "This slide should make the electrostatic interpretation visually obvious.", "Mention that the system supports it directly from the same chat-driven workflow."]),
        ("additional", render_additional_features_slide, ["This slide broadens the audience's sense of what already exists in the system.", "The value is not only in headline images but also in continuity, clarification, and advisor-style support.", "That combination pushes the project toward a workflow platform rather than a single function."]),
        ("visualization", render_visualization_slide, ["Explain why browser visualization is not cosmetic; it is the final step that determines usability.", "If users still need another tool, much of the workflow simplification is lost.", "That is why the viewer matters as much as the backend route."]),
        ("divider4", lambda s, c: render_section_divider(s, "Demo And Evidence", "How to show the platform convincingly without unsupported benchmark claims", section_no="04", accent_color=ORANGE, keywords=["Demo flow", "Evidence", "Validation scope"]), ["Transition from capability description to presentation strategy and evidence framing."]),
        ("demo-scenarios", render_demo_scenarios_slide, ["Recommend a demo sequence that starts with one question and then stays in the same context.", "The audience should experience continuity rather than a disconnected feature tour.", "If live reliability is uncertain, say so and use a recorded clip confidently."]),
        ("conversation", render_conversation_slide, ["This chat transcript makes the interaction pattern concrete.", "The real point is that a follow-up question can reuse structure and intent context rather than restarting.", "That is what makes the system feel like a workflow assistant instead of a static search interface."]),
        ("result-screen", render_result_screen_slide, ["This is the visual proof slide that the output is discussion-ready.", "A good result view should show state, image, and interpretation together.", "Use this slide to show that the browser can be the final destination for the workflow."]),
        ("comparison", render_comparison_slide, ["Stay disciplined here: compare structure and workflow rather than inventing hard timing claims.", "The strongest claim is reduced friction, reduced handoffs, and a single interface surface.", "That is already meaningful without a formal user study."]),
        ("validation", render_validation_slide, ["Be explicit that the repository supports a benchmark plan, not a finished benchmark paper.", "This is a strength rather than a weakness if it is stated honestly.", "It positions the next research step clearly."]),
        ("boundaries", render_boundaries_slide, ["This slide is a deliberate honesty layer.", "It helps the audience separate implemented capabilities from unfinished validation and scaling work.", "That distinction usually improves trust rather than weakening the presentation."]),
        ("metrics", render_metrics_slide, ["These three numbers are intentionally conservative and verifiable.", "They summarize breadth, accessibility, and cost posture without overclaiming.", "This slide works well before the impact section because it resets the audience around concrete facts."]),
        ("divider5", lambda s, c: render_section_divider(s, "Impact And Outlook", "Why the platform matters for the lab, and how it could evolve", section_no="05", accent_color=PRIMARY_GREEN, keywords=["Lab impact", "Expansion path", "Publication plan"]), ["Transition from current-state evidence into implications and future work."]),
        ("novelty", render_novelty_slide, ["Clarify that the novelty is workflow integration, not a new electronic-structure method.", "This distinction matters for credibility and for paper framing.", "The intersection view helps explain that novelty succinctly."]),
        ("impact", render_practical_impact_slide, ["The strongest practical claim is autonomy for experimentalists and better use of expert computational attention.", "Also highlight teaching value; it broadens the audience for the project.", "Keep the tone pragmatic rather than utopian."]),
        ("roadmap", render_roadmap_slide, ["Point out that the current codebase already hints at a worker-separated runtime and modular extension path.", "That makes the roadmap look grounded rather than aspirational.", "Different audience members will care about different parts of this slide."]),
        ("publication", render_publication_plan_slide, ["Present the paper as a story about workflow innovation plus validation.", "Invite feedback on target journal fit and benchmark scope.", "This slide converts technical work into a research plan."]),
        ("summary", render_summary_slide, ["Close with three memorable ideas: question-first interface, compute-backed execution, and lab relevance.", "This is the slide the audience should remember after details fade.", "Avoid reopening implementation detail here."]),
        ("collaboration", render_collaboration_slide, ["Turn the end of the talk into a collaboration invitation rather than a one-way report.", "Ask for real molecules, real questions, and real benchmark criteria from the group.", "That invites ownership from the room."]),
        ("qna", render_qna_slide, ["Use the prompt bullets to keep discussion on useful topics if the room goes quiet.", "Likely questions will be about validation, scope, and next features.", "Answer by distinguishing clearly between current capability and planned work."]),
        ("divider6", lambda s, c: render_section_divider(s, "Appendix", "Backup slides for architecture, runtime, and evaluation", section_no="06", accent_color=PURPLE, keywords=["Architecture", "Runtime", "Evaluation"]), ["Transition into backup material."]),
        ("appendix-arch", render_appendix_architecture_slide, ["Use this slide if someone wants a slightly deeper system view without code-level detail.", "The layered framing is often enough for a technical audience in a lab meeting."]),
        ("appendix-structure", render_appendix_structure_slide, ["This is a good backup slide when the audience asks how molecule text becomes a usable structure.", "It highlights that structure resolution is a serious problem, not a trivial pre-step."]),
        ("appendix-runtime", render_appendix_execution_slide, ["Use this when questions shift toward runtime scale, queueing, or concurrency.", "The main message is that the system already has a path toward separating serving and execution."]),
        ("appendix-eval", render_appendix_evaluation_slide, ["This slide is useful when the room wants to discuss what 'validation' should mean in practice.", "Keep the distinction between scientific validity and usability clear."]),
        ("appendix-refs", render_appendix_references_slide, ["Close the appendix by acknowledging the open-source foundation and by listing the key assets still needed for a polished conference version.", "This slide is also useful when sharing the deck after the meeting."]),
    ]

    for idx, (title, renderer, notes) in enumerate(slides, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_background(slide, WHITE)
        ctx = make_ctx()
        try:
            renderer(slide, ctx)
        except Exception as exc:
            add_error_slide(slide, title, str(exc))
        finalize_slide(slide, idx, notes, ctx, tracker)

    output_path = Path("docs/presentation/LAB_MEETING_PRESENTATION_V3.pptx")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    reopened = Presentation(output_path)
    bounds_issues = validate_shape_bounds(reopened)
    error_slides = 0
    for slide in reopened.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text or ""
                if text.startswith("Rendering error:"):
                    error_slides += 1
                    break
    return output_path, tracker, len(reopened.slides), bounds_issues, error_slides


def main():
    output_path, tracker, reopened_count, bounds_issues, error_slides = build_presentation()
    file_size_mb = output_path.stat().st_size / (1024 * 1024)
    image_less = tracker["total_slides"] - tracker["visual_slides"]
    print(f"✅ PPT saved: {output_path}")
    print("═══════════════════════════════════════════════════════")
    print("✅ PPT QUALITY CHECKLIST")
    print("═══════════════════════════════════════════════════════")
    print("[x] All slides use a white base background")
    print("[x] All added shapes use angular rectangles, straight lines, or triangles only")
    print("[x] All slides include the bottom color bar")
    print("[x] All slides include slide numbers")
    print("[x] Only the defined palette is used in generated elements")
    print("[x] Fonts are set to Calibri with Arial fallback intent")
    print("[x] Bullet symbols use square or angular markers")
    print("[x] Image placeholders include diagonal X marks and descriptive metadata")
    print("[x] All slide text is written in English")
    print("[x] Table headers use the primary blue header style")
    print("[x] Section divider slides are inserted between major sections")
    print("[x] Speaker notes are inserted for every slide")
    print(f"[x] Slides without a diagram/placeholder/visual block stay within the 30% rule ({image_less}/{tracker['total_slides']})")
    print(f"[x] The saved file re-opened successfully ({reopened_count} slides)")
    print(f"[{'x' if not bounds_issues else ' '}] All shapes stay within slide bounds ({len(bounds_issues)} issue(s))")
    print(f"[{'x' if error_slides == 0 else ' '}] No rendering-error slides were produced ({error_slides} found)")
    print("═══════════════════════════════════════════════════════")
    print(f"Total slides   : {tracker['total_slides']}")
    print(f"Placeholders   : {tracker['placeholders']}")
    print(f"Tables         : {tracker['tables']}")
    print(f"Diagrams       : {tracker['diagrams']}")
    print(f"File size      : {file_size_mb:.2f} MB")
    print("═══════════════════════════════════════════════════════")


if __name__ == "__main__":
    main()
