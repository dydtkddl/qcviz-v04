from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


NAVY = RGBColor(0x1B, 0x3A, 0x5C)
GOLD = RGBColor(0xFF, 0xC0, 0x00)
LIGHT = RGBColor(0xF4, 0xF7, 0xFB)
PALE_BLUE = RGBColor(0xE9, 0xF0, 0xF8)
TEXT = RGBColor(0x11, 0x11, 0x11)
MUTED = RGBColor(0x4B, 0x55, 0x63)
BORDER = RGBColor(0xC8, 0xD3, 0xE0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def find_ppt() -> Path:
    matches = list((Path.cwd() / "docs" / "presentation").glob("*withcodex.pptx"))
    if not matches:
        raise FileNotFoundError("Could not find *_withcodex.pptx in docs/presentation")
    return matches[0]


def rgb_from_hex(hex_code: str) -> RGBColor:
    hex_code = hex_code.lstrip("#")
    return RGBColor(int(hex_code[0:2], 16), int(hex_code[2:4], 16), int(hex_code[4:6], 16))


def set_textbox_textbox_style(shape, *, font_name="Segoe UI", font_size=11, color=TEXT, bold=False):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(6)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = font_name
            r.font.size = Pt(font_size)
            r.font.bold = bold
            r.font.color.rgb = color


def add_label(slide, left, top, width, text):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.28))
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY
    box.line.color.rgb = NAVY
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Arial"
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = GOLD
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return box


def add_panel(slide, left, top, width, height, *, fill=LIGHT):
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.line.color.rgb = BORDER
    rect.line.width = Pt(1.1)
    return rect


def add_text(slide, left, top, width, height, text, *, font_name="Segoe UI", font_size=11, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, bullets, *, line_gap=0.52, bullet_color=NAVY, text_color=TEXT, font_size=10.5):
    cur_top = top
    for bullet in bullets:
        marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, cur_top + Inches(0.06), Inches(0.09), Inches(0.09))
        marker.fill.solid()
        marker.fill.fore_color.rgb = bullet_color
        marker.line.color.rgb = bullet_color
        tb = slide.shapes.add_textbox(left + Inches(0.14), cur_top, width - Inches(0.14), Inches(0.42))
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = bullet
        r.font.name = "Segoe UI"
        r.font.size = Pt(font_size)
        r.font.color.rgb = text_color
        cur_top += Inches(line_gap)


def add_tag(slide, left, top, width, text):
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.34))
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = NAVY
    rect.line.width = Pt(1.1)
    tf = rect.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Segoe UI"
    r.font.size = Pt(9.5)
    r.font.bold = True
    r.font.color.rgb = NAVY


def add_strategy_cell(slide, left, top, width, height, title, body, *, fill=WHITE):
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.line.color.rgb = BORDER
    rect.line.width = Pt(1)
    title_box = slide.shapes.add_textbox(left + Inches(0.08), top + Inches(0.06), width - Inches(0.16), Inches(0.28))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Segoe UI"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = NAVY
    body_box = slide.shapes.add_textbox(left + Inches(0.08), top + Inches(0.34), width - Inches(0.16), height - Inches(0.42))
    tf2 = body_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = "Segoe UI"
    r2.font.size = Pt(9.3)
    r2.font.color.rgb = TEXT


def update_slide_22(prs: Presentation) -> None:
    slide = prs.slides[21]

    # Remove any body shapes that may already exist while preserving existing theme header/footer/title.
    keep_names = {
        "Freeform 7",
        "TextBox 9",
        "TextBox 10",
        "Freeform 76",
        "TextBox 77",
        "Slide Number Placeholder 3",
        "TextBox 78",
    }
    for shape in list(slide.shapes):
        if shape.name not in keep_names:
            shape._element.getparent().remove(shape._element)

    # Left impact panel
    add_panel(slide, Inches(0.34), Inches(1.02), Inches(3.0), Inches(3.35), fill=LIGHT)
    add_label(slide, Inches(0.47), Inches(1.14), Inches(1.45), "LAB-LEVEL IMPACT")
    add_text(
        slide,
        Inches(0.48),
        Inches(1.48),
        Inches(2.7),
        Inches(0.62),
        "What changes if this works in a real lab?",
        font_size=16,
        bold=True,
        color=NAVY,
    )
    add_bullets(
        slide,
        Inches(0.52),
        Inches(2.07),
        Inches(2.62),
        [
            "Experimentalists move from request-and-wait workflows to direct computational exploration.",
            "Turnaround shrinks from tool-chain coordination to immediate conversational access.",
            "Computational chemists can spend less time on repetitive setup and more time on validation and interpretation.",
            "Repeated use builds intuition for HOMO, ESP, and geometry optimization inside everyday lab work.",
        ],
        line_gap=0.55,
        font_size=10.2,
    )

    # Right scholarly positioning panel
    add_panel(slide, Inches(3.55), Inches(1.02), Inches(5.92), Inches(1.46), fill=WHITE)
    add_label(slide, Inches(3.68), Inches(1.14), Inches(1.7), "SCHOLARLY POSITIONING")
    add_text(
        slide,
        Inches(3.69),
        Inches(1.47),
        Inches(5.55),
        Inches(0.78),
        "This is not a new electronic-structure theory. It is a platform contribution at the intersection of computational chemistry, scientific workflow automation, and human–AI interaction.",
        font_size=10.6,
        color=TEXT,
    )
    add_tag(slide, Inches(3.72), Inches(2.06), Inches(1.68), "Computational Chemistry")
    add_tag(slide, Inches(5.49), Inches(2.06), Inches(2.02), "Workflow Automation")
    add_tag(slide, Inches(7.62), Inches(2.06), Inches(1.67), "Human–AI Interaction")

    # Right publication strategy panel
    add_panel(slide, Inches(3.55), Inches(2.66), Inches(5.92), Inches(1.71), fill=PALE_BLUE)
    add_label(slide, Inches(3.68), Inches(2.78), Inches(1.62), "PUBLICATION STRATEGY")
    add_strategy_cell(
        slide,
        Inches(3.72),
        Inches(3.12),
        Inches(1.75),
        Inches(0.95),
        "Paper 1",
        "Workflow bottlenecks, platform architecture, and a value/usability demonstration for real lab users.",
        fill=WHITE,
    )
    add_strategy_cell(
        slide,
        Inches(5.63),
        Inches(3.12),
        Inches(1.75),
        Inches(0.95),
        "Paper 2",
        "Expanded workflows, broader engine coverage, and deeper validation across more demanding use cases.",
        fill=WHITE,
    )
    add_strategy_cell(
        slide,
        Inches(7.54),
        Inches(3.12),
        Inches(1.72),
        Inches(0.95),
        "Likely Homes",
        "JCIM, Digital Discovery, and Journal of Cheminformatics as realistic first targets.",
        fill=WHITE,
    )

    # Bottom keynote callout
    callout = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.34), Inches(4.52), Inches(9.13), Inches(0.58))
    callout.fill.solid()
    callout.fill.fore_color.rgb = NAVY
    callout.line.color.rgb = NAVY
    tf = callout.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "The central claim is not that AI replaces chemistry, but that it reduces the friction between research questions and verifiable computation."
    r.font.name = "Segoe UI"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = WHITE


def main() -> None:
    src = find_ppt()
    dst = src.with_name(src.stem + "_slide22done.pptx")
    prs = Presentation(str(src))
    update_slide_22(prs)
    prs.save(str(dst))
    print(dst)


if __name__ == "__main__":
    main()
