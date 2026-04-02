from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


NAVY = RGBColor(0x1B, 0x3A, 0x5C)
GOLD = RGBColor(0xFF, 0xC0, 0x00)
LIGHT = RGBColor(0xF4, 0xF7, 0xFB)
PALE_BLUE = RGBColor(0xE9, 0xF0, 0xF8)
TEXT = RGBColor(0x11, 0x11, 0x11)
BORDER = RGBColor(0xC8, 0xD3, 0xE0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x4B, 0x55, 0x63)


def find_ppt() -> Path:
    matches = list((Path.cwd() / "docs" / "presentation").glob("*slide22done.pptx"))
    if not matches:
        raise FileNotFoundError("Could not find *slide22done.pptx in docs/presentation")
    return matches[0]


def add_label(slide, left, top, width, text):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.28))
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY
    box.line.color.rgb = NAVY
    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Arial"
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = GOLD
    return box


def add_panel(slide, left, top, width, height, fill):
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.line.color.rgb = BORDER
    rect.line.width = Pt(1.0)
    return rect


def add_text(slide, left, top, width, height, text, *, font_name="Segoe UI", font_size=11, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, bullets, *, line_gap=0.37, font_size=9.8):
    current = top
    for bullet in bullets:
        marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, current + Inches(0.07), Inches(0.085), Inches(0.085))
        marker.fill.solid()
        marker.fill.fore_color.rgb = NAVY
        marker.line.color.rgb = NAVY
        tb = slide.shapes.add_textbox(left + Inches(0.13), current, width - Inches(0.13), Inches(0.34))
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = bullet
        r.font.name = "Segoe UI"
        r.font.size = Pt(font_size)
        r.font.color.rgb = TEXT
        current += Inches(line_gap)


def add_roadmap_cell(slide, left, top, width, height, title, body):
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = BORDER
    rect.line.width = Pt(1.0)
    add_text(slide, left + Inches(0.08), top + Inches(0.06), width - Inches(0.16), Inches(0.23), title, font_size=9.8, color=NAVY, bold=True)
    add_text(slide, left + Inches(0.08), top + Inches(0.30), width - Inches(0.16), height - Inches(0.36), body, font_size=9.0, color=TEXT)


def update_slide_23(prs: Presentation) -> None:
    slide = prs.slides[22]

    keep_names = {
        "Freeform 7",
        "TextBox 9",
        "TextBox 10",
        "Freeform 76",
        "TextBox 77",
        "Slide Number Placeholder 3",
        "TextBox 78",
    }
    for shape in list(slide.shapes):
        if shape.name not in keep_names:
            shape._element.getparent().remove(shape._element)

    # Intro summary band
    intro = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.34), Inches(0.95), Inches(9.12), Inches(0.70))
    intro.fill.solid()
    intro.fill.fore_color.rgb = NAVY
    intro.line.color.rgb = NAVY
    tf = intro.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "The current platform already covers isolated-molecule tasks such as single-point energy, orbital visualisation, ESP, and geometry optimisation — but the architecture is designed to extend well beyond that starting point."
    r.font.name = "Segoe UI"
    r.font.size = Pt(10.5)
    r.font.color.rgb = WHITE
    r.font.bold = True

    # Left column: deeper workflows
    add_panel(slide, Inches(0.34), Inches(1.82), Inches(4.45), Inches(2.55), LIGHT)
    add_label(slide, Inches(0.47), Inches(1.94), Inches(2.02), "1. DEEPER MOLECULAR WORKFLOWS")
    add_text(slide, Inches(0.50), Inches(2.28), Inches(4.1), Inches(0.46), "Where do we go next at the molecule level?", font_size=15.5, color=NAVY, bold=True)
    add_bullets(
        slide,
        Inches(0.53),
        Inches(2.76),
        Inches(4.02),
        [
            "Transition-state search and IRC as connected multi-step reaction workflows.",
            "Frequency analysis for IR/Raman spectra and thermochemical data.",
            "NMR chemical-shift prediction as a more complete interpretation layer.",
            "A move from one-off calculations toward end-to-end reaction-analysis workflows.",
        ],
        line_gap=0.43,
        font_size=10.0,
    )

    # Right upper: more complex systems
    add_panel(slide, Inches(4.96), Inches(1.82), Inches(4.50), Inches(1.44), WHITE)
    add_label(slide, Inches(5.09), Inches(1.94), Inches(2.00), "2. MORE COMPLEX PHYSICAL SYSTEMS")
    add_bullets(
        slide,
        Inches(5.12),
        Inches(2.31),
        Inches(4.05),
        [
            "From cluster DFT to surface DFT and ultimately periodic DFT.",
            "An expansion from isolated molecules toward materials-scale research questions.",
            "A broader system scope without discarding the conversational front-end.",
        ],
        line_gap=0.39,
        font_size=9.8,
    )

    # Right middle: feasibility
    add_panel(slide, Inches(4.96), Inches(3.42), Inches(4.50), Inches(0.95), PALE_BLUE)
    add_label(slide, Inches(5.09), Inches(3.54), Inches(1.55), "WHY THIS IS FEASIBLE")
    add_bullets(
        slide,
        Inches(5.12),
        Inches(3.88),
        Inches(4.06),
        [
            "Intent interpretation is already separated from engine execution.",
            "New backends and workflows can be added as modular extensions.",
        ],
        line_gap=0.31,
        font_size=9.4,
    )

    # Bottom roadmap
    add_panel(slide, Inches(0.34), Inches(4.50), Inches(9.12), Inches(0.62), WHITE)
    add_label(slide, Inches(0.47), Inches(4.62), Inches(1.62), "PUBLICATION ROADMAP")
    add_roadmap_cell(
        slide,
        Inches(2.25),
        Inches(4.58),
        Inches(2.18),
        Inches(0.42),
        "Paper 1",
        "Platform architecture and molecule-level workflow.",
    )
    add_roadmap_cell(
        slide,
        Inches(4.56),
        Inches(4.58),
        Inches(2.18),
        Inches(0.42),
        "Paper 2+",
        "Multi-step workflows and spectroscopy expansion.",
    )
    add_roadmap_cell(
        slide,
        Inches(6.87),
        Inches(4.58),
        Inches(2.18),
        Inches(0.42),
        "Later Scope",
        "Periodic systems and expanded validation.",
    )


def main() -> None:
    src = find_ppt()
    dst = src.with_name(src.stem.replace("slide22done", "slide23done") + src.suffix)
    prs = Presentation(str(src))
    update_slide_23(prs)
    prs.save(str(dst))
    print(dst)


if __name__ == "__main__":
    main()
